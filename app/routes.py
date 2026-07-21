from flask import Blueprint, render_template, redirect, url_for, flash, request, jsonify
from flask_login import login_user, logout_user, login_required, current_user
from werkzeug.utils import secure_filename
from app import db
from app.models import User, Notebook, Document, ChatMessage, Quiz, Question, QuizAttempt, Answer
from app.ai_helper import answer_question, generate_quiz
from datetime import datetime
import os
from PyPDF2 import PdfReader

main = Blueprint('main', __name__)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'txt', 'pdf', 'ppt', 'pptx'}

# ==================== AUTHENTICATION ROUTES ====================

@main.route('/')
def index():
    if current_user.is_authenticated:
        return redirect(url_for('main.home'))
    return redirect(url_for('main.login'))

@main.route('/register', methods=['GET', 'POST'])
def register():
    if current_user.is_authenticated:
        return redirect(url_for('main.home'))
    
    if request.method == 'POST':
        name = request.form.get('name')
        email = request.form.get('email')
        password = request.form.get('password')
        
        if User.query.filter_by(email=email).first():
            flash('Email already registered!', 'danger')
            return redirect(url_for('main.register'))
        
        user = User(name=name, email=email)
        user.set_password(password)
        db.session.add(user)
        db.session.commit()
        
        flash('Registration successful! Please login.', 'success')
        return redirect(url_for('main.login'))
    
    return render_template('register.html')

@main.route('/login', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        return redirect(url_for('main.home'))
    
    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')
        
        user = User.query.filter_by(email=email).first()
        
        if user and user.check_password(password):
            login_user(user)
            flash('Login successful!', 'success')
            return redirect(url_for('main.home'))
        else:
            flash('Invalid email or password!', 'danger')
    
    return render_template('login.html')

@main.route('/logout')
@login_required
def logout():
    logout_user()
    flash('Logged out successfully!', 'info')
    return redirect(url_for('main.login'))

# ==================== MAIN APP ROUTES ====================

@main.route('/home')
@login_required
def home():
    notebooks = Notebook.query.filter_by(user_id=current_user.id).order_by(Notebook.created_at.desc()).all()
    return render_template('home.html', notebooks=notebooks)

# ==================== NOTEBOOK ROUTES ====================

@main.route('/notebook/create', methods=['POST'])
@login_required
def create_notebook():
    name = request.form.get('name')
    description = request.form.get('description', '')
    
    if not name:
        flash('Notebook name is required!', 'danger')
        return redirect(url_for('main.home'))
    
    notebook = Notebook(
        user_id=current_user.id,
        name=name,
        description=description
    )
    db.session.add(notebook)
    db.session.commit()
    
    flash(f'Notebook "{name}" created!', 'success')
    return redirect(url_for('main.view_notebook', notebook_id=notebook.id))

@main.route('/notebook/<int:notebook_id>')
@login_required
def view_notebook(notebook_id):
    notebook = Notebook.query.get_or_404(notebook_id)
    
    if notebook.user_id != current_user.id:
        flash('Access denied!', 'danger')
        return redirect(url_for('main.home'))
    
    chats = ChatMessage.query.filter_by(notebook_id=notebook_id).order_by(ChatMessage.created_at.asc()).all()
    
    return render_template('notebook.html', notebook=notebook, chats=chats)

@main.route('/notebook/<int:notebook_id>/delete', methods=['POST'])
@login_required
def delete_notebook(notebook_id):
    notebook = Notebook.query.get_or_404(notebook_id)
    
    if notebook.user_id != current_user.id:
        flash('Access denied!', 'danger')
        return redirect(url_for('main.home'))
    
    for doc in notebook.documents:
        if doc.filename:
            filepath = os.path.join('app/static/uploads', doc.filename)
            if os.path.exists(filepath):
                os.remove(filepath)
    
    db.session.delete(notebook)
    db.session.commit()
    
    flash('Notebook deleted!', 'success')
    return redirect(url_for('main.home'))

# ==================== DOCUMENT ROUTES ====================

@main.route('/notebook/<int:notebook_id>/upload', methods=['POST'])
@login_required
def upload_document(notebook_id):
    notebook = Notebook.query.get_or_404(notebook_id)
    
    if notebook.user_id != current_user.id:
        flash('Access denied!', 'danger')
        return redirect(url_for('main.home'))
    
    if 'file' not in request.files:
        flash('No file selected!', 'danger')
        return redirect(url_for('main.view_notebook', notebook_id=notebook_id))
    
    file = request.files['file']
    
    if file.filename == '':
        flash('No file selected!', 'danger')
        return redirect(url_for('main.view_notebook', notebook_id=notebook_id))
    
    if file and allowed_file(file.filename):
        original_filename = secure_filename(file.filename)
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{current_user.id}_{timestamp}_{original_filename}"
        
        filepath = os.path.join('app/static/uploads', filename)
        file.save(filepath)
        
        content = ''
        file_type = original_filename.rsplit('.', 1)[1].lower()
        
        try:
            if file_type == 'txt':
                with open(filepath, 'r', encoding='utf-8') as f:
                    content = f.read()
            
            elif file_type == 'pdf':
                reader = PdfReader(filepath)
                for page in reader.pages:
                    content += page.extract_text() + "\n"
            
            elif file_type in ['ppt', 'pptx']:
                from pptx import Presentation
                prs = Presentation(filepath)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
        
        except Exception as e:
            flash(f'Error reading file: {str(e)}', 'danger')
            os.remove(filepath)
            return redirect(url_for('main.view_notebook', notebook_id=notebook_id))
        
        if len(content.strip()) < 50:
            flash('File content is too short or empty!', 'warning')
            os.remove(filepath)
            return redirect(url_for('main.view_notebook', notebook_id=notebook_id))
        
        document = Document(
            notebook_id=notebook_id,
            filename=filename,
            original_filename=original_filename,
            content=content,
            file_type=file_type
        )
        db.session.add(document)
        db.session.commit()
        
        flash(f'Document "{original_filename}" uploaded successfully!', 'success')
    else:
        flash('Invalid file type! Only .txt, .pdf, .ppt, and .pptx allowed.', 'danger')
    
    return redirect(url_for('main.view_notebook', notebook_id=notebook_id))

@main.route('/document/<int:document_id>/delete', methods=['POST'])
@login_required
def delete_document(document_id):
    document = Document.query.get_or_404(document_id)
    notebook = document.notebook
    
    if notebook.user_id != current_user.id:
        flash('Access denied!', 'danger')
        return redirect(url_for('main.home'))
    
    if document.filename:
        filepath = os.path.join('app/static/uploads', document.filename)
        if os.path.exists(filepath):
            os.remove(filepath)
    
    db.session.delete(document)
    db.session.commit()
    
    flash('Document deleted!', 'success')
    return redirect(url_for('main.view_notebook', notebook_id=notebook.id))

# ==================== CHAT ROUTES ====================

@main.route('/notebook/<int:notebook_id>/chat', methods=['POST'])
@login_required
def chat(notebook_id):
    notebook = Notebook.query.get_or_404(notebook_id)
    
    if notebook.user_id != current_user.id:
        return jsonify({'error': 'Access denied'}), 403
    
    data = request.get_json()
    if not data:
        return jsonify({'error': 'No data provided'}), 400
    
    user_message = data.get('message', '').strip()
    
    if not user_message:
        return jsonify({'error': 'Empty message'}), 400
    
    documents = Document.query.filter_by(notebook_id=notebook_id).all()
    
    if not documents:
        ai_response = "Please upload some documents first so I can help you study! 📚"
    else:
        if 'quiz' in user_message.lower() or 'test' in user_message.lower() or 'questions' in user_message.lower():
            import re
            num_match = re.search(r'\d+', user_message)
            num_questions = int(num_match.group()) if num_match else 5
            num_questions = min(max(num_questions, 3), 10)
            
            ai_response = f"I'll generate a quiz with {num_questions} questions for you! Click the 'Generate Quiz' button in the sidebar to start. 🎯"
        else:
            try:
                ai_response = answer_question(user_message, documents)
            except Exception as e:
                print(f"Error in answer_question: {e}")
                ai_response = "I'm having trouble processing your question right now. Please try again! 🤔"
    
    chat_message = ChatMessage(
        notebook_id=notebook_id,
        user_message=user_message,
        ai_response=ai_response
    )
    db.session.add(chat_message)
    db.session.commit()
    
    return jsonify({
        'success': True,
        'user_message': user_message,
        'ai_response': ai_response,
        'timestamp': chat_message.created_at.strftime('%I:%M %p')
    })

# ==================== QUIZ ROUTES ====================

@main.route('/notebook/<int:notebook_id>/generate-quiz', methods=['POST'])
@login_required
def generate_quiz_route(notebook_id):
    notebook = Notebook.query.get_or_404(notebook_id)
    
    if notebook.user_id != current_user.id:
        flash('Access denied!', 'danger')
        return redirect(url_for('main.home'))
    
    documents = Document.query.filter_by(notebook_id=notebook_id).all()
    
    if not documents:
        flash('Please upload documents first!', 'warning')
        return redirect(url_for('main.view_notebook', notebook_id=notebook_id))
    
    num_questions = int(request.form.get('num_questions', 5))
    num_questions = min(max(num_questions, 3), 10)
    
    questions_data = generate_quiz(documents, num_questions)
    
    if not questions_data:
        flash('Failed to generate quiz. Please try again.', 'danger')
        return redirect(url_for('main.view_notebook', notebook_id=notebook_id))
    
    quiz = Quiz(
        notebook_id=notebook_id,
        title=f"Quiz - {notebook.name} ({datetime.now().strftime('%b %d, %Y')})"
    )
    db.session.add(quiz)
    db.session.flush()
    
    for q_data in questions_data:
        question = Question(
            quiz_id=quiz.id,
            question_text=q_data['text'],
            option_a=q_data['option_a'],
            option_b=q_data['option_b'],
            option_c=q_data['option_c'],
            option_d=q_data['option_d'],
            correct_answer=q_data['correct_answer']
        )
        db.session.add(question)
    
    db.session.commit()
    
    flash(f'Quiz generated with {len(questions_data)} questions!', 'success')
    return redirect(url_for('main.take_quiz', quiz_id=quiz.id))

@main.route('/quiz/<int:quiz_id>')
@login_required
def take_quiz(quiz_id):
    quiz = Quiz.query.get_or_404(quiz_id)
    
    if quiz.notebook.user_id != current_user.id:
        flash('Access denied!', 'danger')
        return redirect(url_for('main.home'))
    
    questions = Question.query.filter_by(quiz_id=quiz_id).all()
    
    return render_template('quiz.html', quiz=quiz, questions=questions)

@main.route('/quiz/<int:quiz_id>/submit', methods=['POST'])
@login_required
def submit_quiz(quiz_id):
    quiz = Quiz.query.get_or_404(quiz_id)
    
    if quiz.notebook.user_id != current_user.id:
        flash('Access denied!', 'danger')
        return redirect(url_for('main.home'))
    
    questions = Question.query.filter_by(quiz_id=quiz_id).all()
    
    correct_count = 0
    total_questions = len(questions)
    
    percentage = 0
    attempt = QuizAttempt(
        quiz_id=quiz_id,
        user_id=current_user.id,
        score=0,
        total_questions=total_questions,
        percentage=0
    )
    db.session.add(attempt)
    db.session.flush()
    
    for question in questions:
        user_answer = request.form.get(f'question_{question.id}', '')
        is_correct = (user_answer == question.correct_answer)
        
        if is_correct:
            correct_count += 1
        
        answer = Answer(
            attempt_id=attempt.id,
            question_id=question.id,
            selected_answer=user_answer,
            is_correct=is_correct
        )
        db.session.add(answer)
    
    percentage = (correct_count / total_questions * 100) if total_questions > 0 else 0
    attempt.score = correct_count
    attempt.percentage = percentage
    
    db.session.commit()
    
    flash(f'Quiz submitted! Score: {correct_count}/{total_questions} ({percentage:.0f}%)', 'success')
    return redirect(url_for('main.quiz_results', attempt_id=attempt.id))

@main.route('/quiz/results/<int:attempt_id>')
@login_required
def quiz_results(attempt_id):
    attempt = QuizAttempt.query.get_or_404(attempt_id)
    
    if attempt.user_id != current_user.id:
        flash('Access denied!', 'danger')
        return redirect(url_for('main.home'))
    
    quiz = attempt.quiz
    questions = Question.query.filter_by(quiz_id=quiz.id).all()
    
    return render_template('results.html', attempt=attempt, quiz=quiz, questions=questions)

# Add this route to your routes.py file

@main.route('/performance')
@login_required
def performance():
    """
    Quiz Performance History Dashboard
    """
    # Get all quiz attempts for current user
    attempts = QuizAttempt.query.filter_by(user_id=current_user.id).order_by(QuizAttempt.attempted_at.desc()).all()
    
    # Calculate statistics
    total_quizzes = len(attempts)
    
    if total_quizzes > 0:
        avg_score = sum(a.percentage for a in attempts) / total_quizzes
        highest_score = max(a.percentage for a in attempts)
        lowest_score = min(a.percentage for a in attempts)
        
        # Count performance levels
        excellent = sum(1 for a in attempts if a.percentage >= 80)
        good = sum(1 for a in attempts if 60 <= a.percentage < 80)
        fair = sum(1 for a in attempts if 40 <= a.percentage < 60)
        poor = sum(1 for a in attempts if a.percentage < 40)
        
        # Recent trend (last 5 quizzes)
        recent_attempts = attempts[:5]
        recent_scores = [a.percentage for a in reversed(recent_attempts)]
        
        # Subject performance (group by notebook)
        subject_stats = {}
        for attempt in attempts:
            notebook_name = attempt.quiz.notebook.name
            if notebook_name not in subject_stats:
                subject_stats[notebook_name] = {
                    'total': 0,
                    'scores': []
                }
            subject_stats[notebook_name]['total'] += 1
            subject_stats[notebook_name]['scores'].append(attempt.percentage)
        
        # Calculate average per subject
        for subject in subject_stats:
            scores = subject_stats[subject]['scores']
            subject_stats[subject]['average'] = sum(scores) / len(scores)
    else:
        avg_score = 0
        highest_score = 0
        lowest_score = 0
        excellent = good = fair = poor = 0
        recent_scores = []
        subject_stats = {}
    
    return render_template('performance.html', 
                         attempts=attempts,
                         total_quizzes=total_quizzes,
                         avg_score=avg_score,
                         highest_score=highest_score,
                         lowest_score=lowest_score,
                         excellent=excellent,
                         good=good,
                         fair=fair,
                         poor=poor,
                         recent_scores=recent_scores,
                         subject_stats=subject_stats)